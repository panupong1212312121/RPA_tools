from img2table.document import Image
import tabula
import fitz as pyMuPDF
from docx import Document
from pptx import Presentation

import pandas as pd

from transformer import Transformer
from widget import *

import requests

class ExtractTable:
    
    def image(self,file_afterTransfrom):
        dfs = Image(file_afterTransfrom)
        return dfs
    
    def pdf(self,file,file_afterTransfrom):
        dfs = {}
        page_amount = pyMuPDF.open(file_afterTransfrom).page_count
        for page_num in range(page_amount):
            tables = tabula.read_pdf(file, pages=page_num+1)
            if tables:
                dfs[page_num+1] = [table for table in tables]
        return dfs
    
    def docx(self,file_afterTransfrom):
        document = Document(file_afterTransfrom)
        dfs = []
        for table in document.tables:
            df = [['' for _ in range(len(table.columns))] for _ in range(len(table.rows))]
            for i, row in enumerate(table.rows):
                for j, cell in enumerate(row.cells):
                    if cell.text:
                        df[i][j] = cell.text
            dfs.append(pd.DataFrame(df))
        return dfs
    
    def pptx(self,file_afterTransfrom):
        presentation = Presentation(file_afterTransfrom)
        dfs = {}
        page = 0
        for slide in presentation.slides:
            df = []
            page += 1
            hasTable = True
            for shape in slide.shapes:
                if not shape.has_table:
                    hasTable = False
                    continue
                table = shape.table
                row_count = len(table.rows)
                col_count = len(table.columns)
                table_data = []
                for row in range(row_count):
                    row_values = []
                    for col in range(col_count):
                        cell = table.cell(row, col)
                        text = ""
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text += run.text
                        row_values.append(text)
                    table_data.append(row_values)
                df.append(pd.DataFrame(table_data))
            if hasTable:
                dfs[page] = df
        return dfs
    
    def link(self,currentState):
        dfs = {}
        for key in currentState:
            if 'link' in key and 'amount' not in key:
                try:
                    url_name = currentState[key]
                    df = pd.read_html(url_name)
                    dfs[key] = df
                # except requests.exceptions.HTTPError as err:
                except:
                    print(f'{url_name} is invalid name')
        return dfs
        
    def run(self,uploaded_files,currentState,all_null):
        transform_obj = Transformer()
        ############## Loading Pages ################



        ################################################
        if len(uploaded_files)>0:
            for file in uploaded_files:

                filename_input = file.name
                file_extension = file.type.split('/')[-1].lower()

                file_afterTransfrom = transform_obj.fileUtil(file,filename_input)

                if file_extension in ['png', 'jpg','jpeg']:
                    data = self.image(file_afterTransfrom)
                elif file_extension in ['pdf']:
                    data = self.pdf(file,file_afterTransfrom)
                elif file_extension in ['vnd.openxmlformats-officedocument.wordprocessingml.document']:
                    data = self.docx(file_afterTransfrom)
                elif file_extension in ['vnd.openxmlformats-officedocument.presentationml.presentation']:
                    data = self.pptx(file_afterTransfrom)

                transform_obj.toExcel(data,file_extension,filename_input)

        if len(currentState)>2 and all_null==False:
            data = self.link(currentState)

            transform_obj.toExcel(data)

            
            
        