from img2table.document import Image
import tabula
import fitz as pyMuPDF
from docx import Document
from pptx import Presentation

import pandas as pd

from algo.transformer import Transformer
from ui.table_extractor import *

import streamlit as st

from ui.table_extractor import FileUploader,SaveButton

class ExtractTable:
    
    def image(self,file_afterTransfrom):
        dfs = Image(file_afterTransfrom)
        return dfs
    
    def pdf(self,file,file_afterTransfrom):
        dfs = {}
        page_amount = pyMuPDF.open(file_afterTransfrom).page_count
        for page_num in range(page_amount):
            tables = tabula.read_pdf(file, pages=page_num+1)
            if tables:
                dfs[page_num+1] = [table for table in tables]
        return dfs
    
    def docx(self,file_afterTransfrom):
        document = Document(file_afterTransfrom)
        dfs = []
        for table in document.tables:
            df = [['' for _ in range(len(table.columns))] for _ in range(len(table.rows))]
            for i, row in enumerate(table.rows):
                for j, cell in enumerate(row.cells):
                    if cell.text:
                        df[i][j] = cell.text
            dfs.append(pd.DataFrame(df))
        return dfs
    
    def pptx(self,file_afterTransfrom):
        presentation = Presentation(file_afterTransfrom)
        dfs = {}
        page = 0
        for slide in presentation.slides:
            df = []
            page += 1
            hasTable = True
            for shape in slide.shapes:
                if not shape.has_table:
                    hasTable = False
                    continue
                table = shape.table
                row_count = len(table.rows)
                col_count = len(table.columns)
                table_data = []
                for row in range(row_count):
                    row_values = []
                    for col in range(col_count):
                        cell = table.cell(row, col)
                        text = ""
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text += run.text
                        row_values.append(text)
                    table_data.append(row_values)
                df.append(pd.DataFrame(table_data))
            if hasTable:
                dfs[page] = df
        return dfs
    
    def link(self,currentState):
        dfs = {}
        all_url_list_name = {}
        invalid_url_list_name = []
        for key in currentState:
            if 'url' in key and 'amount' not in key:
                try:
                    url_name = currentState[key]
                    df = pd.read_html(url_name)
                    dfs[key] = df
                    all_url_list_name[key] = url_name
                except:
                    invalid_url_list_name.append(url_name)
        return dfs,all_url_list_name,invalid_url_list_name
        
    def run(self,currentState):
        transform_obj = Transformer()
        file_uploader_obj = FileUploader()
        save_button_obj = SaveButton()

        file_list = st.session_state.get(file_uploader_obj.key_file_list)
        has_link = st.session_state.get(save_button_obj.key_criteria_data_save_button)

        if len(file_list)>0:
            files_list_no_detected_table = []
            urls_list_no_detected_table = []

            for file in file_list:

                filename_input = file.name
                file_extension = file.type.split('/')[-1].lower()

                file_afterTransfrom = transform_obj.fileUtil(file,filename_input)

                if file_extension in ['png', 'jpg','jpeg']:
                    data = self.image(file_afterTransfrom)
                elif file_extension in ['pdf']:
                    data = self.pdf(file,file_afterTransfrom)
                elif file_extension in ['vnd.openxmlformats-officedocument.wordprocessingml.document']:
                    data = self.docx(file_afterTransfrom)
                elif file_extension in ['vnd.openxmlformats-officedocument.presentationml.presentation']:
                    data = self.pptx(file_afterTransfrom)

                files_no_detected_table,urls_no_detected_table = transform_obj.toExcel(data,
                                                                                       file_extension,
                                                                                       filename_input)

                files_list_no_detected_table.extend(files_no_detected_table)
                urls_list_no_detected_table.extend(urls_no_detected_table)

                transform_obj.deleteNewFile(filename_input)

            transform_obj.toTxt(files_list_no_detected_table,
                                urls_list_no_detected_table)

        if len(currentState)>5 and has_link:
            dfs,all_url_list_name,invalid_url_list_name = self.link(currentState)

            files_list_no_detected_table,urls_list_no_detected_table = transform_obj.toExcel(dfs,
                                                                                             all_url_name=all_url_list_name)
            transform_obj.toTxt(files_list_no_detected_table,
                                urls_list_no_detected_table,
                                invalid_url_list_name)

            
            
        